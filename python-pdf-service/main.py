import fitz  # PyMuPDF
import time
import re
from fastapi import FastAPI, File, UploadFile, HTTPException
import uvicorn
import os
import tempfile
import shutil
import threading
from pptx import Presentation
from mistralai import Mistral

# Enable output for debugging
os.environ["PYTHONUNBUFFERED"] = "1"

print("🐍 Starting Python PDF Service...")
print("📦 Importing dependencies...")

try:
    print("✅ PyMuPDF imported successfully")
    print("✅ FastAPI imported successfully")
    print("✅ All dependencies loaded")
except Exception as e:
    print(f"❌ Import error: {e}")
    raise

app = FastAPI(title="PDF Text Extraction Service", version="1.0.0")
print("🚀 FastAPI app created successfully")

# Initialize Mistral client
mistral_client = None
try:
    api_key = os.environ.get("MISTRAL_API_KEY")
    if api_key:
        mistral_client = Mistral(api_key=api_key)
        print("✅ Mistral client initialized successfully")
    else:
        print("⚠️ MISTRAL_API_KEY not found - PPTX OCR functionality will be disabled")
except Exception as e:
    print(f"⚠️ Failed to initialize Mistral client: {e}")


@app.get("/health")
async def health_check():
    """Health check endpoint"""
    print("🏥 Health check requested")
    return {"status": "healthy", "service": "pdf-extraction"}


@app.post("/extract-text")
async def extract_pdf_text(file: UploadFile = File(...)):
    """
    Extract text from PDF using PyMuPDF
    Returns structured text with page information and metadata
    """
    start_time = time.time()

    try:
        # Validate file type
        if not file.filename.lower().endswith('.pdf'):
            raise HTTPException(
                status_code=400, detail="Only PDF files are supported")

        # Read file content
        pdf_content = await file.read()

        # Open PDF with PyMuPDF - maximum speed settings
        pdf_document = fitz.open(stream=pdf_content, filetype="pdf")

        # Ultra-fast text extraction - minimal object creation
        pages_text = []
        total_chars = 0
        page_count = pdf_document.page_count

        # Optimized loop - pre-allocate and minimal operations
        for page_num in range(page_count):
            page_text = pdf_document[page_num].get_text()
            char_count = len(page_text)

            pages_text.append({
                "page_number": page_num + 1,
                "text": page_text,
                "char_count": char_count
            })
            total_chars += char_count

        # Minimal metadata - only what's needed
        metadata = {
            "title": pdf_document.metadata.get("title", ""),
            "author": pdf_document.metadata.get("author", ""),
            "subject": pdf_document.metadata.get("subject", ""),
            "creator": pdf_document.metadata.get("creator", ""),
            "producer": pdf_document.metadata.get("producer", ""),
            "creation_date": pdf_document.metadata.get("creationDate", ""),
            "modification_date": pdf_document.metadata.get("modDate", "")
        }

        # Close document
        pdf_document.close()

        processing_time = time.time() - start_time

        # Optimized result construction - minimal operations
        result = {
            "success": True,
            "filename": file.filename,
            "pages": pages_text,
            "metadata": {
                "total_pages": page_count,
                "total_characters": total_chars,
                **metadata
            },
            "processing_time_seconds": round(processing_time, 3),
            "extraction_method": "PyMuPDF"
        }

        # No logging for maximum speed
        return result

    except Exception as e:
        # Minimal error handling for speed
        raise HTTPException(
            status_code=500, detail=f"PDF processing failed: {str(e)}")


def extract_images_from_pptx(pptx_path, temp_dir):
    """Extract images from PPTX file and return list of image paths"""
    try:
        prs = Presentation(pptx_path)
        image_paths = []
        image_count = 0

        for slide_num, slide in enumerate(prs.slides, 1):
            for shape in slide.shapes:
                if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                    try:
                        image = shape.image
                        ext = image.ext
                        image_count += 1
                        image_filename = f"slide_{slide_num}_image_{image_count}.{ext}"
                        image_path = os.path.join(temp_dir, image_filename)

                        with open(image_path, "wb") as f:
                            f.write(image.blob)

                        image_paths.append(image_path)
                    except Exception as e:
                        print(
                            f"⚠️ Failed to extract image from slide {slide_num}: {e}")
                        continue

        return image_paths
    except Exception as e:
        raise Exception(f"Failed to extract images from PPTX: {e}")


def clean_ocr_text(text):
    """Clean and format OCR extracted text"""
    if not text:
        return ""

    # Remove image references like ![img-0.jpeg](img-0.jpeg)
    text = re.sub(r'!\[img-\d+\.\w+\]\(img-\d+\.\w+\)', '', text)

    # Remove markdown headers - convert to plain text
    text = re.sub(r'^#{1,6}\s*', '', text,
                  flags=re.MULTILINE)  # Remove # headers

    # Remove markdown bold/italic formatting
    text = re.sub(r'\*\*([^*]+)\*\*', r'\1', text)  # Remove **bold**
    text = re.sub(r'\*([^*]+)\*', r'\1', text)      # Remove *italic*
    text = re.sub(r'__([^_]+)__', r'\1', text)      # Remove __bold__
    text = re.sub(r'_([^_]+)_', r'\1', text)        # Remove _italic_

    # Convert markdown bullet points to simple bullets
    text = re.sub(r'^\s*[-*+]\s+', '• ', text, flags=re.MULTILINE)

    # Clean up excessive whitespace
    text = re.sub(r'\s+', ' ', text)  # Multiple spaces to single space
    # Multiple newlines to double
    text = re.sub(r'\n\s*\n\s*\n+', '\n\n', text)

    # Clean up line breaks and remove empty lines, but preserve structure
    lines = text.split('\n')
    cleaned_lines = []
    prev_was_empty = False

    for line in lines:
        stripped_line = line.strip()

        # Keep empty lines for structure, but not more than 1 consecutive
        if not stripped_line:
            if not prev_was_empty:
                cleaned_lines.append('')
            prev_was_empty = True
        else:
            cleaned_lines.append(stripped_line)
            prev_was_empty = False

    # Join and final cleanup
    result = '\n'.join(cleaned_lines).strip()

    # Final cleanup - ensure no excessive spacing
    result = re.sub(r'\n\n+', '\n\n', result)  # Max 2 consecutive newlines

    return result


def ocr_image_with_mistral(image_path, client):
    """Process a single image with Mistral OCR"""
    try:
        start = time.time()

        with open(image_path, "rb") as f:
            uploaded_img = client.files.upload(
                file={
                    "file_name": os.path.basename(image_path),
                    "content": f,
                },
                purpose="ocr"
            )

        ocr_response = client.ocr.process(
            model="mistral-ocr-latest",
            document={
                "type": "file",
                "file_id": uploaded_img.id
            },
            include_image_base64=False
        )

        elapsed = time.time() - start

        # Extract text from the OCR response
        ocr_text = ""
        if hasattr(ocr_response, 'pages') and ocr_response.pages:
            # Extract markdown content from OCR pages
            for page in ocr_response.pages:
                if hasattr(page, 'markdown') and page.markdown:
                    ocr_text += page.markdown + "\n"
        elif hasattr(ocr_response, 'text'):
            ocr_text = ocr_response.text
        elif hasattr(ocr_response, 'content'):
            if isinstance(ocr_response.content, list):
                ocr_text = "\n".join([item.text if hasattr(item, 'text') else str(
                    item) for item in ocr_response.content])
            else:
                ocr_text = str(ocr_response.content)
        else:
            ocr_text = str(ocr_response)

        # Clean up the extracted text
        ocr_text = clean_ocr_text(ocr_text.strip())

        return {
            "image_path": os.path.basename(image_path),
            "text": ocr_text,
            "processing_time": elapsed,
            "char_count": len(ocr_text)
        }

    except Exception as e:
        return {
            "image_path": os.path.basename(image_path),
            "text": f"[OCR Error: {str(e)}]",
            "processing_time": 0,
            "char_count": 0
        }


def process_images_parallel(image_paths, client):
    """Process multiple images with Mistral OCR in parallel"""
    if not client:
        raise Exception("Mistral client not initialized")

    results = [None] * len(image_paths)
    threads = []

    def ocr_worker(idx, image_path):
        results[idx] = ocr_image_with_mistral(image_path, client)

    # Start all threads
    for idx, image_path in enumerate(image_paths):
        thread = threading.Thread(target=ocr_worker, args=(idx, image_path))
        threads.append(thread)
        thread.start()

    # Wait for all threads to complete
    for thread in threads:
        thread.join()

    return [r for r in results if r is not None]


@app.post("/process-pptx")
async def process_pptx(file: UploadFile = File(...)):
    """
    Process PPTX file by extracting images and running OCR on them using Mistral
    Returns structured text with image OCR results and metadata
    """
    start_time = time.time()
    temp_dir = None

    try:
        # Validate file type
        if not file.filename.lower().endswith('.pptx'):
            raise HTTPException(
                status_code=400, detail="Only PPTX files are supported")

        if not mistral_client:
            raise HTTPException(
                status_code=500, detail="Mistral OCR service not available - MISTRAL_API_KEY not configured")

        # Create temporary directory for processing
        temp_dir = tempfile.mkdtemp()

        # Save uploaded file
        pptx_path = os.path.join(temp_dir, file.filename)
        with open(pptx_path, "wb") as buffer:
            shutil.copyfileobj(file.file, buffer)

        # Extract images from PPTX
        print(f"🖼️ Extracting images from {file.filename}...")
        image_paths = extract_images_from_pptx(pptx_path, temp_dir)

        if not image_paths:
            # No images found, return minimal response
            return {
                "success": True,
                "filename": file.filename,
                "pages": [{
                    "page_number": 1,
                    "text": "[No images found in PPTX file]",
                    "char_count": 0
                }],
                "metadata": {
                    "total_pages": 1,
                    "total_characters": 0,
                    "images_processed": 0,
                    "pdf_conversion_time": 0,
                    "ocr_processing_time": 0,
                    "source_type": "pptx_no_images"
                },
                "processing_time_seconds": round(time.time() - start_time, 3),
                "extraction_method": "pptx_image_ocr"
            }

        print(f"📊 Found {len(image_paths)} images, starting OCR processing...")

        # Process all images with OCR in parallel
        ocr_start = time.time()
        ocr_results = process_images_parallel(image_paths, mistral_client)
        ocr_time = time.time() - ocr_start

        # Combine all OCR results into pages
        pages = []
        total_chars = 0

        print("\n📄 PPTX Processing Results:")
        print("=" * 60)

        for idx, result in enumerate(ocr_results):
            # Format the page text with better structure
            cleaned_text = result['text']
            if cleaned_text and not cleaned_text.startswith('[OCR Error'):
                # Just use the cleaned text without image reference
                page_text = cleaned_text
            else:
                # Keep error messages as is
                page_text = cleaned_text

            pages.append({
                "page_number": idx + 1,
                "text": page_text,
                "char_count": len(page_text)
            })
            total_chars += len(page_text)

            # Print each OCR result
            print(f"\n🖼️ Image {idx + 1}: {result['image_path']}")
            print(f"⏱️ Processing time: {result['processing_time']:.2f}s")
            print(f"📝 Characters extracted: {result['char_count']}")
            print(f"📄 Full OCR Text:\n{result['text']}")
            print("-" * 40)

        processing_time = time.time() - start_time

        result = {
            "success": True,
            "filename": file.filename,
            "pages": pages,
            "metadata": {
                "total_pages": len(pages),
                "total_characters": total_chars,
                "images_processed": len(image_paths),
                "pdf_conversion_time": 0,  # Not applicable for PPTX
                "ocr_processing_time": round(ocr_time, 3),
                "source_type": "pptx_images"
            },
            "processing_time_seconds": round(processing_time, 3),
            "extraction_method": "pptx_image_ocr"
        }

        print(
            f"\n✅ PPTX processing complete: {len(pages)} images processed, {total_chars} characters extracted")
        print(f"🚀 Total processing time: {processing_time:.2f} seconds")
        print("=" * 60)

        return result

    except Exception as e:
        raise HTTPException(
            status_code=500, detail=f"PPTX processing failed: {str(e)}")

    finally:
        # Cleanup temporary directory
        if temp_dir and os.path.exists(temp_dir):
            try:
                shutil.rmtree(temp_dir)
            except Exception as e:
                print(f"⚠️ Failed to cleanup temp directory: {e}")


@app.get("/")
async def root():
    """Root endpoint with service information"""
    return {
        "service": "PDF Text Extraction Service",
        "version": "1.0.0",
        "description": "High-performance PDF text extraction using PyMuPDF and PPTX processing with Mistral OCR",
        "endpoints": {
            "health": "/health",
            "extract_pdf": "/extract-text",
            "process_pptx": "/process-pptx"
        },
        "features": {
            "pdf_extraction": "PyMuPDF",
            "pptx_processing": "Image extraction + Mistral OCR",
            "mistral_ocr_available": mistral_client is not None
        }
    }

if __name__ == "__main__":
    print("🌟 Starting uvicorn server on 0.0.0.0:8000...")
    try:
        uvicorn.run(
            app,
            host="0.0.0.0",
            port=8000,
            log_level="info",      # Enable logging for debugging
            access_log=True,       # Enable access logs for debugging
            workers=1              # Single worker for speed
        )
    except Exception as e:
        print(f"❌ Server startup failed: {e}")
        raise
